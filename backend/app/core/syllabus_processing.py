"""Document text extraction using PyMuPDF (PDFs) and python-pptx (presentations).

PyMuPDF is used instead of pypdf because it:
  - Handles complex layouts, multi-column text, and embedded fonts correctly
  - Preserves reading order across columns
  - Extracts tables and structured content reliably
  - Is ~10× faster for large documents

OCR fallback (Tesseract via pytesseract + pdf2image / Pillow):
  - Activated automatically when PyMuPDF returns < 100 chars/page (scanned PDFs)
  - Supports direct image uploads (PNG, JPG, JPEG, WEBP, TIFF, BMP)
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import List

logger = logging.getLogger(__name__)

# ── Characters per page threshold below which we consider a page "scanned" ──
_OCR_THRESHOLD = 100


def extract_text_from_file(path: str) -> str:
    """Extract clean plain-text from PDF, PPTX/PPT, TXT, or image files."""
    file_path = Path(path)
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(str(file_path))
    if suffix in (".pptx", ".ppt"):
        return _extract_pptx(str(file_path))
    if suffix in (".txt", ".md"):
        return file_path.read_text(errors="replace")
    if suffix in (".png", ".jpg", ".jpeg", ".webp", ".tiff", ".tif", ".bmp"):
        return _extract_image(str(file_path))

    return ""


def extract_pages_from_file(path: str) -> dict[int, str]:
    """Return a dict mapping 1-based page number → text for that page/slide.

    PDFs: one entry per page.
    PPTX: one entry per slide.
    Other formats: single entry {1: full_text}.
    """
    file_path = Path(path)
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        try:
            import fitz
            doc = fitz.open(str(file_path))
            pages: dict[int, str] = {}
            for i, page in enumerate(doc):
                text = page.get_text("text").strip()
                if text:
                    pages[i + 1] = text
            doc.close()
            return pages
        except Exception:
            return {1: extract_text_from_file(path)}

    if suffix in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            slides: dict[int, str] = {}
            for idx, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.extend(para.text.strip() for para in shape.text_frame.paragraphs if para.text.strip())
                if parts:
                    slides[idx] = "\n".join(parts)
            return slides
        except Exception:
            return {1: extract_text_from_file(path)}

    # Default: full text as single "page"
    return {1: extract_text_from_file(path)}


def find_topic_pages(topics: list[str], page_texts: dict[int, str]) -> dict[str, int]:
    """Return {topic_name: first_page_number} for each topic found in the document.

    Matching is case-insensitive and checks whether a substantial part of the
    topic name (≥ 4 contiguous words or the full name if shorter) appears on
    the page.
    """
    topic_pages: dict[str, int] = {}
    for topic in topics:
        topic_low = topic.lower().strip()
        if len(topic_low) < 3:
            continue
        # For topics longer than ~30 chars, match on first 4 words (more robust)
        words = topic_low.split()
        search_str = " ".join(words[:4]) if len(words) > 4 else topic_low
        for page_num, text in sorted(page_texts.items()):
            if search_str in text.lower():
                topic_pages[topic] = page_num
                break
    return topic_pages


def _extract_pdf(path: str) -> str:
    """Use PyMuPDF for high-quality PDF text extraction, with OCR fallback
    for scanned/image-only pages."""
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(path)
        pages: list[str] = []
        ocr_page_indices: list[int] = []

        for i, page in enumerate(doc):
            text = page.get_text("text", flags=fitz.TEXT_PRESERVE_LIGATURES | fitz.TEXT_MEDIABOX_CLIP)
            stripped = text.strip()
            if stripped:
                pages.append(stripped)
                if len(stripped) < _OCR_THRESHOLD:
                    # Page likely has embedded text but it may be very sparse;
                    # keep but flag for potential supplemental OCR.
                    ocr_page_indices.append(i)
            else:
                # Completely blank extraction — almost certainly a scanned page.
                ocr_page_indices.append(i)
                pages.append("")   # placeholder, filled by OCR below

        doc.close()

        # Perform OCR on pages that had insufficient text
        if ocr_page_indices:
            ocr_texts = _ocr_pdf_pages(path, ocr_page_indices)
            for idx, ocr_text in zip(ocr_page_indices, ocr_texts):
                # Replace blank placeholder or supplement sparse text
                if idx < len(pages) and (not pages[idx] or len(pages[idx]) < _OCR_THRESHOLD):
                    pages[idx] = ocr_text

        return "\n\n".join(p for p in pages if p.strip())

    except Exception as exc:
        logger.warning("PyMuPDF extraction failed for %s: %s", path, exc)
        # Last resort: try full OCR on the whole PDF
        return _ocr_pdf_pages_all(path)


def _ocr_pdf_pages(path: str, page_indices: list[int]) -> list[str]:
    """Convert specific PDF pages to images and run Tesseract OCR."""
    try:
        from pdf2image import convert_from_path
        import pytesseract

        # Convert only the required pages (1-indexed in pdf2image)
        first = min(page_indices) + 1
        last  = max(page_indices) + 1
        images = convert_from_path(path, first_page=first, last_page=last, dpi=200)

        results: list[str] = []
        img_iter = iter(images)
        for i, page_idx in enumerate(range(first - 1, last)):
            img = next(img_iter, None)
            if img is None:
                results.append("")
                continue
            if page_idx in page_indices:
                text = pytesseract.image_to_string(img, lang="eng")
                results.append(text.strip())
        return results

    except ImportError:
        logger.debug("pdf2image/pytesseract not installed — OCR skipped")
        return [""] * len(page_indices)
    except Exception as exc:
        logger.warning("OCR failed for %s: %s", path, exc)
        return [""] * len(page_indices)


def _ocr_pdf_pages_all(path: str) -> str:
    """Run OCR on all pages of a PDF (last-resort fallback)."""
    try:
        from pdf2image import convert_from_path
        import pytesseract

        images = convert_from_path(path, dpi=200)
        return "\n\n".join(
            pytesseract.image_to_string(img, lang="eng").strip()
            for img in images
        )
    except ImportError:
        logger.debug("pdf2image/pytesseract not installed — OCR skipped")
        return ""
    except Exception as exc:
        logger.warning("Full OCR fallback failed for %s: %s", path, exc)
        return ""


def _extract_image(path: str) -> str:
    """Run Tesseract OCR on a standalone image file (PNG, JPG, TIFF, etc.)."""
    try:
        from PIL import Image
        import pytesseract

        img = Image.open(path)
        text = pytesseract.image_to_string(img, lang="eng")
        return text.strip()
    except ImportError:
        logger.debug("Pillow/pytesseract not installed — image OCR skipped")
        return ""
    except Exception as exc:
        logger.warning("Image OCR failed for %s: %s", path, exc)
        return ""


def _extract_pptx(path: str) -> str:
    """Extract text from every slide, including speaker notes."""
    try:
        from pptx import Presentation

        prs = Presentation(path)
        slides: list[str] = []
        for idx, slide in enumerate(prs.slides, 1):
            parts: list[str] = [f"--- Slide {idx} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            parts.append(line)
            # Include speaker notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                notes_text = notes_tf.text.strip() if notes_tf else ""
                if notes_text:
                    parts.append(f"[Notes] {notes_text}")
            slides.append("\n".join(parts))
        return "\n\n".join(slides)
    except Exception as exc:
        logger.warning("pptx extraction failed for %s: %s", path, exc)
        return ""


# ── Topic / heading extraction ────────────────────────────────────────────────

# Patterns that look like section headings or topic lines
_HEADING_RE = re.compile(
    r"^(\d+[\.\)]\s+|Chapter|Unit|Module|Section|Topic|Lecture|Lab|Week)\s*",
    re.IGNORECASE,
)
_MIN_TOPIC_LEN = 4
_MAX_TOPIC_LEN = 120

# Phrases that indicate institutional boilerplate rather than academic topics
_JUNK_PHRASES: frozenset = frozenset({
    "institute of", "college of", "university of", "department of",
    "school of", "faculty of", "b.tech", "m.tech", "b.e.", "m.e.",
    "b.sc", "m.sc",
    "title of the course", "name of the course", "course title",
    "course code", "l t p c", "l  t  p  c", "credits",
    "professional elective", "open elective", "honours",
    "co-requisite", "pre-requisite", "pre requisite",
    "contact hours", "examination pattern", "bloom",
    "taught by", "offered by",
    "text book", "reference book", "bibliography",
    "regulation", "programme outcome", "course outcome",
    "total:", "grand total", "page ",
})
# Regulation code suffixes common in Indian university syllabi (R17, R20, R22 …)
_REG_CODE_RE = re.compile(r"\br(1[5-9]|2\d)\b", re.IGNORECASE)
# Fraction of uppercase alpha chars above which a line is considered a header
_ALLCAPS_FRACTION = 0.78


def _is_junk_line(line: str) -> bool:
    """Return True if *line* is institutional boilerplate, not an academic topic."""
    low = line.lower()
    if any(phrase in low for phrase in _JUNK_PHRASES):
        return True
    if _REG_CODE_RE.search(line):
        return True
    alpha = [c for c in line if c.isalpha()]
    if len(alpha) >= 5 and sum(c.isupper() for c in alpha) / len(alpha) > _ALLCAPS_FRACTION:
        return True
    # Pure numbers / separators with no real words
    if re.fullmatch(r"[\d\s\.\-\|\/:,;()\[\]]+", line.strip()):
        return True
    return False


def split_into_topics(text: str, max_topics: int = 200) -> List[str]:
    """Extract clean academic topic names from raw document text.

    Heading-like lines are preferred. Institutional boilerplate
    (college names, ALL-CAPS headers, form fields) is filtered out.
    Results are deduplicated.
    """
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    seen: set = set()

    def _unique_add(lst: list, line: str) -> None:
        key = re.sub(r"\s+", " ", line.lower())
        if key not in seen:
            seen.add(key)
            lst.append(line)

    # First pass: heading-like lines that pass the junk filter
    headings: list = []
    for line in lines:
        if (
            _HEADING_RE.search(line)
            and _MIN_TOPIC_LEN <= len(line) <= _MAX_TOPIC_LEN
            and not _is_junk_line(line)
        ):
            _unique_add(headings, line)

    if len(headings) >= 5:
        return headings[:max_topics]

    # Second pass: all non-trivial, non-junk lines
    topics: list = []
    for line in lines:
        if (
            _MIN_TOPIC_LEN <= len(line) <= _MAX_TOPIC_LEN
            and not _is_junk_line(line)
            and not re.fullmatch(r"[\d\s\.\-\|\/\\]+", line.strip())
        ):
            _unique_add(topics, line)

    return topics[:max_topics]
